"""Regenerate docs/speaker-notes.md from the speaker notes in the deck.

The .pptx is the single source of truth: those notes are what the
presenter actually sees on the night. Run this after editing any note so
the rehearsal markdown cannot drift out of sync:

    python docs/sync_speaker_notes.py

Note format expected in each slide's notes pane:

    [~45 s]

    <the text to say>

    (Cue: <stage direction>)
"""

import re
from pathlib import Path

from pptx import Presentation

DOCS = Path(__file__).resolve().parent
DECK = DOCS / "presentation.pptx"
OUT = DOCS / "speaker-notes.md"

HEADER = """# Speaker script

> The full text to say, slide by slide (~9-10 minutes total), plus
> stage cues. Rehearse from this, then deliver it in your own
> phrasing — reading aloud verbatim always sounds worse than it reads.
>
> Generated from presentation.pptx by sync_speaker_notes.py — edit the
> notes in the deck, then re-run, so the two never disagree.

## Speaker assignment (balanced by speaking time)

| Speaker | Slides | ~time |
|---|---|---|
| Izabella Atajanyan | 1–4 (title, rules, why, formulation) | ~3.3 min |
| Ani Baghumyan | 5–8 (solvers, system, setup, inference results) | ~4.0 min |
| Anna Gasparyan | 9–13 (local search, difficulty, robustness, conclusions, Q&A lead) | ~3.6 min |

Swap freely — just keep the three shares roughly equal (equal input is
graded) and hand over cleanly between slides 4→5 and 8→9.
"""


def slide_title(slide):
    """First non-trivial text on the slide (skips clue digits, page nos).

    Titles may wrap across lines inside their text frame, so join them
    rather than taking only the first line.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if len(text) > 3:
                return " ".join(text.split())
    return "(untitled)"


def parse_note(raw):
    """Split a note into (timing, spoken text, [extra notes]).

    Notes are parsed by blank-line-separated blocks rather than by one
    regex, because a slide may carry several parenthetical blocks (a
    stage cue *and* a Q&A note) and the spoken text itself may run to
    several paragraphs.
    """
    timing, spoken, extras = "", [], []

    for block in re.split(r"\n\s*\n", raw.strip()):
        block = block.strip()
        if not block:
            continue
        m = re.fullmatch(r"\[([^\]]*)\]", block)
        if m and not timing:
            timing = m.group(1).strip()
        elif block.startswith("("):
            extras.append(block.strip("()").strip())
        else:
            spoken.append(block)

    return timing, "\n\n".join(spoken), extras


def main():
    pres = Presentation(DECK)
    parts = [HEADER]

    for i, slide in enumerate(pres.slides, 1):
        raw = (slide.notes_slide.notes_text_frame.text
               if slide.has_notes_slide else "")
        timing, say, extras = parse_note(raw)
        heading = f"## Slide {i} — {slide_title(slide)}"
        if timing:
            heading += f"  ({timing})"
        parts.append(heading + "\n")
        parts.append((f"**Say:** {say}\n" if say else "*(nothing spoken)*\n"))
        for extra in extras:
            parts.append(f"> {extra}\n")

    OUT.write_text("\n".join(parts), encoding="utf-8")
    print(f"wrote {OUT} from {len(pres.slides)} slides")


if __name__ == "__main__":
    main()
